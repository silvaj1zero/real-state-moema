# -*- coding: utf-8 -*-
"""PITCH DE CAPTAÇÃO v1 (Rua Alvorada, 116 — apto 52, Ed. Manuela) — .pptx editável.
Design RE/MAX Galeria (mesmo DS do pitch Honduras). 16 slides, 16:9.
Fonte dos números: ACM v2 jun/2026 (squad acm-imobiliario) — AP 51 comparável de ouro,
anúncio R$ 1,55–1,70 mi, fechamento R$ 1,35–1,50 mi, piso R$ 1,25 mi, venal R$ 1,10 mi."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

HERE = os.path.dirname(os.path.abspath(__file__))

# paleta RE/MAX
AZUL=RGBColor(0x00,0x3D,0xA5); AZUL900=RGBColor(0x01,0x2a,0x6e); AZUL700=RGBColor(0x01,0x3a,0x8f)
VERM=RGBColor(0xDC,0x14,0x31); VERMC=RGBColor(0xff,0x5a,0x73); BRANCO=RGBColor(0xFF,0xFF,0xFF)
GELO=RGBColor(0xf4,0xf6,0xfb); CINZA100=RGBColor(0xee,0xf1,0xf6); CINZA300=RGBColor(0xcf,0xd6,0xe4)
CINZA500=RGBColor(0x7d,0x8a,0xa3); CINZA700=RGBColor(0x3a,0x42,0x58); TINTA=RGBColor(0x1a,0x22,0x33)
ICE=RGBColor(0x9f,0xc0,0xff); ICE2=RGBColor(0xcf,0xe0,0xff); OURO=RGBColor(0xE8,0xA3,0x17); LAR=RGBColor(0xC8,0x65,0x1F)
OK=RGBColor(0x1f,0x9d,0x57)
HEAD="Inter"; BODY="Inter"
SW, SH = 13.333, 7.5
MX = 0.62
CW = SW - 2*MX

prs = Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
BLANK = prs.slide_layouts[6]

def _noshadow(sp): sp.shadow.inherit=False
def rect(s, x,y,w,h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp=s.shapes.add_shape(shape, Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line: sp.line.color.rgb=line; sp.line.width=Pt(0.8)
    else: sp.line.fill.background()
    _noshadow(sp); return sp

def bg(s, color=BRANCO): rect(s,0,0,SW,SH,color)
def brandbar(s):
    rect(s,0,0,SW*0.70,0.10,AZUL); rect(s,SW*0.70,0,SW*0.30,0.10,VERM)

def txt(s, x,y,w,h, runs, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
        italic=False, anchor=MSO_ANCHOR.TOP, ls=1.0, space=4):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    lines = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=ls; p.space_after=Pt(space)
        segs = line if isinstance(line, list) else [(line, color, bold)]
        for seg in segs:
            t, c, bd = (seg+ (None,)*3)[:3]
            r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.name=font
            r.font.bold=bold if bd is None else bd; r.font.italic=italic
            r.font.color.rgb=color if c is None else c
    return b

def foot(s, n, dark=False):
    c = BRANCO if dark else AZUL
    b=s.shapes.add_textbox(Inches(MX),Inches(SH-0.5),Inches(6),Inches(0.3)); tf=b.text_frame
    tf.margin_left=0;tf.margin_top=0; p=tf.paragraphs[0]
    for t,col in [("RE/",c),("MAX",VERM if not dark else VERMC),(" Galeria",c)]:
        r=p.add_run(); r.text=t; r.font.size=Pt(10.5); r.font.bold=True; r.font.name=HEAD; r.font.color.rgb=col
    pg=s.shapes.add_textbox(Inches(SW-1.4),Inches(SH-0.5),Inches(0.8),Inches(0.3)); ptf=pg.text_frame
    ptf.margin_left=0;ptf.margin_right=0; pp=ptf.paragraphs[0]; pp.alignment=PP_ALIGN.RIGHT
    rr=pp.add_run(); rr.text=str(n); rr.font.size=Pt(11); rr.font.name=BODY
    rr.font.color.rgb=(ICE2 if dark else CINZA500)

def bug(s, x, y):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(4),Inches(0.7)); tf=b.text_frame
    tf.margin_left=0;tf.margin_top=0; p=tf.paragraphs[0]
    for t,c in [("RE/",BRANCO),("MAX",VERMC)]:
        r=p.add_run(); r.text=t; r.font.size=Pt(26); r.font.bold=True; r.font.name=HEAD; r.font.color.rgb=c
    p2=tf.add_paragraph(); r=p2.add_run(); r.text="GALERIA · MOEMA"; r.font.size=Pt(11); r.font.bold=True
    r.font.name=BODY; r.font.color.rgb=ICE

def kicker_title(s, kicker, title, tsize=30):
    txt(s, MX, 0.55, CW, 0.4, kicker.upper(), 13, VERM, bold=True)
    txt(s, MX, 0.95, CW, 1.1, title, tsize, AZUL900, bold=True, font=HEAD, ls=1.05)

def card(s, x,y,w,h, accent=False, gelo=True, fill=None, line=None):
    f = fill if fill is not None else (AZUL if accent else (GELO if gelo else BRANCO))
    sp=rect(s,x,y,w,h, f, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            line=(line if line is not None else (None if accent else CINZA300)))
    sp.adjustments[0]=0.06; return sp

def stat(s, x,y,w,h, num, unit, label, accent=False, numsize=34):
    card(s,x,y,w,h,accent=accent)
    numc = BRANCO if accent else AZUL
    lblc = ICE2 if accent else CINZA700
    b=s.shapes.add_textbox(Inches(x+0.28),Inches(y+0.22),Inches(w-0.56),Inches(h-0.4)); tf=b.text_frame
    tf.word_wrap=True; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]
    r=p.add_run(); r.text=num; r.font.size=Pt(numsize); r.font.bold=True; r.font.name=HEAD; r.font.color.rgb=numc
    if unit:
        ru=p.add_run(); ru.text=unit; ru.font.size=Pt(17); ru.font.bold=True; ru.font.name=HEAD; ru.font.color.rgb=(ICE2 if accent else CINZA500)
    p2=tf.add_paragraph(); p2.space_before=Pt(6); p2.line_spacing=1.05
    r2=p2.add_run(); r2.text=label; r2.font.size=Pt(12.5); r2.font.name=BODY; r2.font.color.rgb=lblc

def bullets(s, x,y,w, items, size=15, gap=8, color=CINZA700):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(4.5)); tf=b.text_frame
    tf.word_wrap=True; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.line_spacing=1.15
        rb=p.add_run(); rb.text="▪  "; rb.font.size=Pt(size); rb.font.name=BODY; rb.font.color.rgb=VERM; rb.font.bold=True
        segs = it if isinstance(it, list) else [(it, color, None)]
        for seg in segs:
            t,c,bd=(seg+(None,)*3)[:3]
            r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.name=BODY
            r.font.color.rgb=color if c is None else c; r.font.bold=False if bd is None else bd
    return b

def checklist(s, x,y,w, items, size=14, gap=9):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(5)); tf=b.text_frame
    tf.word_wrap=True; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.line_spacing=1.12
        rc=p.add_run(); rc.text="✔  "; rc.font.size=Pt(size); rc.font.name=BODY; rc.font.color.rgb=OK; rc.font.bold=True
        r=p.add_run(); r.text=it; r.font.size=Pt(size); r.font.name=BODY; r.font.color.rgb=CINZA700
    return b

def sourcebox(s, x,y,w,h, runs, size=12.5):
    sp=rect(s,x,y,w,h,GELO, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=CINZA300); sp.adjustments[0]=0.04
    rect(s,x,y,0.06,h,AZUL)
    txt(s,x+0.22,y+0.16,w-0.4,h-0.3, runs, size, CINZA700, anchor=MSO_ANCHOR.MIDDLE, ls=1.15)

def table(s, x, y, w, headers, rows, colw, hl_row=None, avg_row=None, fsize=11.5, rowh=0.34):
    """tabela estilo .acm: header azul900, linha destacada azul, linha média gelo."""
    n_rows = len(rows)+1
    shp = s.shapes.add_table(n_rows, len(headers), Inches(x), Inches(y), Inches(w), Inches(rowh*n_rows))
    tb = shp.table
    tb.first_row = False; tb.horz_banding = False
    for j, cwd in enumerate(colw): tb.columns[j].width = Inches(cwd)
    for j, htxt in enumerate(headers):
        cell = tb.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = AZUL900
        cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = htxt; r.font.size=Pt(fsize-1); r.font.bold=True; r.font.name=BODY; r.font.color.rgb=BRANCO
    for i, row in enumerate(rows, start=1):
        is_hl = (hl_row is not None and i-1 == hl_row)
        is_avg = (avg_row is not None and i-1 == avg_row)
        for j, val in enumerate(row):
            cell = tb.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = AZUL if is_hl else (GELO if is_avg else BRANCO)
            cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(val); r.font.size=Pt(fsize); r.font.name=BODY
            r.font.bold = is_hl or is_avg
            r.font.color.rgb = BRANCO if is_hl else (TINTA if is_avg else CINZA700)
    return shp

def newslide(dark=False):
    s=prs.slides.add_slide(BLANK); bg(s, AZUL900 if dark else BRANCO)
    if not dark: brandbar(s)
    return s

# ============================================================ SLIDES (16)

# S1 CAPA
s=prs.slides.add_slide(BLANK); bg(s,AZUL900); brandbar(s)
txt(s,MX,1.05,CW,0.4,"RE/MAX GALERIA · PROPOSTA DE REPRESENTAÇÃO EXCLUSIVA",14,ICE,bold=True)
txt(s,MX,1.55,8.6,2.4,"Saudações, Dorothy e família. Vamos vender o seu imóvel pelo melhor valor real.",38,BRANCO,bold=True,font=HEAD,ls=1.05)
txt(s,MX,4.25,7.6,0.9,"RE/MAX, Real Estate Maximum: o máximo de serviço ao cliente.",17,ICE2,ls=1.3)
txt(s,MX,5.2,8.5,0.4,"Apto 52 · Ed. Manuela — Rua Alvorada, 116 · Vila Olímpia, São Paulo · Agosto de 2026",15,BRANCO,bold=True)
txt(s,7.7,4.5,5.0,2.2,[
    [("Luciana Borba",BRANCO,True)],
    [("Consultora Imobiliária · RE/MAX Galeria · (11) 99995-2014",ICE2,False)],
    [("CRECI: 045063-J",ICE,False)],
    [(" ",ICE,False)],
    [("Rodrigo Lima · (11) 96371-0014 · CRECI: 242.934",ICE2,False)],
    [("Alameda dos Jurupis, 452 · Conj. 153 B · Moema",ICE2,False)],
],12.5,ICE2,align=PP_ALIGN.RIGHT,ls=1.25,space=2)
bug(s,MX,6.2)

# S2 RE/MAX EM NÚMEROS (+ DNA)
s=newslide(); kicker_title(s,"Quem confia a venda à RE/MAX está com a maior rede do mundo","A força que vai trabalhar pelo seu imóvel",tsize=28)
txt(s,MX,1.95,CW,0.8,[[("Fundada em 1973, a RE/MAX é a ",CINZA700,False),("maior rede imobiliária do mundo em volume de vendas",TINTA,True),(" — e também a número 1 no Brasil.",CINZA700,False)]],16,CINZA700,ls=1.35)
sp=rect(s,MX,2.78,CW,0.92,GELO,shape=MSO_SHAPE.ROUNDED_RECTANGLE,line=CINZA300); sp.adjustments[0]=0.05
rect(s,MX,2.78,0.07,0.92,VERM)
txt(s,MX+0.28,2.9,CW-0.5,0.7,[[("Faz parte do DNA da RE/MAX",AZUL900,True),(" proporcionar ao cliente uma experiência com 100% de ética, comprometimento, excelência e resultado — revolucionando o mercado com método eficaz e inovador.",CINZA700,False)]],13.5,CINZA700,italic=True,ls=1.25,anchor=MSO_ANCHOR.MIDDLE)
cw3=(CW-0.6)/3; y=4.0; hh=2.1
stat(s,MX,y,cw3,hh,"115","+ países","144 mil+ corretores conectados pelo mesmo sistema, há mais de 50 anos")
stat(s,MX+cw3+0.3,y,cw3,hh,"600","+ lojas","No Brasil, em mais de 250 cidades, com 11 mil+ profissionais")
stat(s,MX+2*(cw3+0.3),y,cw3,hh,"R$ 14,3","bi","Em VGV no Brasil em 2025 — número 1 em vendas no país",accent=True)
foot(s,2)

# S3 CAPACIDADE DE VENDA
s=newslide(); kicker_title(s,"Capacidade de venda · rede + método","Vendemos em menos da metade do tempo")
cw2=(CW-1.2)/2; y=2.15; hh=2.0
card(s,MX,y,cw2,hh,gelo=True)
txt(s,MX,y+0.35,cw2,1.0,"490",48,CINZA700,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
txt(s,MX+0.3,y+1.35,cw2-0.6,0.5,"dias — prazo médio do mercado tradicional (OLX/Zap)",13,CINZA700,align=PP_ALIGN.CENTER,ls=1.15)
txt(s,MX+cw2,y+0.55,1.2,0.8,"→",28,VERM,bold=True,align=PP_ALIGN.CENTER)
xr=MX+cw2+1.2
card(s,xr,y,cw2,hh,accent=True)
txt(s,xr,y+0.35,cw2,1.0,"140",48,BRANCO,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
txt(s,xr+0.3,y+1.35,cw2-0.6,0.5,"dias — prazo médio de vendas da RE/MAX no Brasil",13,BRANCO,align=PP_ALIGN.CENTER,ls=1.15)
yb=4.45; hb=2.1
card(s,MX,yb,3.7,hb,accent=True)
txt(s,MX+0.35,yb+0.3,3.0,1.0,[[("7",BRANCO,True),(" de 10",ICE2,True)]],46,BRANCO,bold=True,font=HEAD)
txt(s,MX+0.35,yb+1.25,3.0,0.7,"imóveis vendidos pela RE/MAX têm o comprador originado por outra imobiliária ou corretor parceiro.",12.5,ICE2,ls=1.25)
txt(s,MX+4.1,yb+0.2,CW-(4.1),hb,[
    [("Seu imóvel fica disponível para ",CINZA700,False),("milhares de corretores",TINTA,True),(" — mas com ",CINZA700,False),("uma só interlocutora",TINTA,True),(": a Luciana.",CINZA700,False)],
    [(" ",CINZA700,False)],
    [("A rede amplia o alcance; a exclusividade garante a consistência. ",CINZA700,False),("Mais alcance, mesma mensagem.",AZUL,True)],
],17,CINZA700,ls=1.35,space=6,anchor=MSO_ANCHOR.MIDDLE)
foot(s,3)

# S4 PLANO DE COMERCIALIZAÇÃO
s=newslide(); kicker_title(s,"Plano de comercialização · como a RE/MAX vende","Uma venda boa e rápida nasce de estratégia, planejamento e execução",tsize=26)
txt(s,MX,2.0,CW,0.5,[[("Na RE/MAX, cada imóvel recebe um ",CINZA700,False),("plano de comercialização estruturado",TINTA,True),(":",CINZA700,False)]],16,CINZA700,ls=1.25)
cw2=(CW-0.5)/2
checklist(s,MX,2.65,cw2,[
    "Uma só corretora como interlocutora",
    "Precificação baseada em análise de mercado",
    "Produção de fotos e materiais de alta qualidade",
    "Divulgação direcionada nos principais canais",
    "Compartilhamento com ampla rede de corretores",
],size=14,gap=10)
checklist(s,MX+cw2+0.5,2.65,cw2,[
    "Qualificação de compradores",
    "Acompanhamento das visitas",
    "Gestão das propostas",
    "Condução da negociação",
    "Relatórios periódicos de desempenho e feedback (a combinar)",
],size=14,gap=10)
sp=rect(s,MX,5.7,CW,0.85,AZUL,shape=MSO_SHAPE.ROUNDED_RECTANGLE); sp.adjustments[0]=0.10
txt(s,MX+0.3,5.78,CW-0.6,0.7,[[("O resultado: ",BRANCO,True),("mais visibilidade, mais visitas qualificadas e negociações mais assertivas.",ICE2,False)]],16,ICE2,anchor=MSO_ANCHOR.MIDDLE,ls=1.2)
foot(s,4)

# S5 DIVISOR ACM
s=newslide(dark=True)
txt(s,MX,1.6,CW,1.6,"ACM",90,RGBColor(0x2a,0x4d,0x9e),bold=True,font=HEAD)
txt(s,MX,3.15,10.5,1.2,"O seu imóvel e o valor que o mercado realmente paga",36,BRANCO,bold=True,font=HEAD,ls=1.05)
txt(s,MX,4.6,9.5,1.0,"A Análise Comparativa de Mercado — o método mais eficiente no mundo e no Brasil — apoiada em vendas registradas, não em achismo.",17,ICE2,ls=1.4)
foot(s,5,dark=True)

# S6 SOBRE O IMÓVEL
s=newslide(); kicker_title(s,"Sobre o seu imóvel · características confirmadas","3 dormitórios com 2 vagas num prédio boutique de só 20 unidades",tsize=26)
txt(s,MX,2.0,7.2,1.0,[[("No coração da Vila Olímpia, o ",CINZA700,False),("Ed. Manuela (1994)",TINTA,True),(" tem apenas ",CINZA700,False),("2 apartamentos por andar",TINTA,True),(". O apto 52 fica no ",CINZA700,False),("5º de 10 andares",TINTA,True),(", de frente para a rua, ",CINZA700,False),("pronto para morar",TINTA,True),(".",CINZA700,False)]],14.5,CINZA700,ls=1.3)
bullets(s,MX,3.15,7.2,[
    [("3 dormitórios (1 suíte)",TINTA,True),(" + lavabo + varanda + dependência completa.",CINZA700,False)],
    [("Conservação boa:",TINTA,True),(" pintado, piso de madeira restaurado, ventilação cruzada.",CINZA700,False)],
    [("Lazer completo",TINTA,True),(" — academia, sauna e salão de festas — rateado entre poucas famílias.",CINZA700,False)],
],size=14,gap=9)
sourcebox(s,MX,5.35,7.2,1.15,[[("Due diligence antes do anúncio: ",TINTA,True),("confirmar o carnê de IPTU do próprio apto 52 (a área oficial veio do apto 71, mesmo prédio) e a averbação das 2 vagas na matrícula.",CINZA700,False)]],size=11.5)
xr=8.2; cw2s=(SW-MX-xr-0.25)/2
stat(s,xr,2.0,cw2s,2.05,"113"," m²","Área útil",accent=True)
stat(s,xr+cw2s+0.25,2.0,cw2s,2.05,"213"," m²","Área cadastral (IPTU)")
stat(s,xr,4.25,cw2s,2.05,"2","","Vagas privativas")
stat(s,xr+cw2s+0.25,4.25,cw2s,2.05,"20","","Unidades no prédio (2 por andar)")
foot(s,6)

# S7 DUAS RÉGUAS
s=newslide(); kicker_title(s,"Antes dos números · a distinção mais importante","Duas réguas de área — uma para calcular, outra para vender",tsize=27)
cw2=(CW-0.5)/2; y=2.1; hh=2.5
card(s,MX,y,cw2,hh,gelo=True)
txt(s,MX+0.3,y+0.25,cw2-0.6,0.7,[[("213 m² ",AZUL,True),("cadastrais",CINZA500,True)]],26,AZUL,bold=True,font=HEAD)
txt(s,MX+0.3,y+0.95,cw2-0.6,1.4,[[("A régua de CÁLCULO. ",TINTA,True),("É a área do IPTU, do ITBI e da matrícula — inclui a privativa, a parte proporcional das áreas comuns e as 2 vagas. É com ela que comparamos venda com venda.",CINZA700,False)]],12.5,CINZA700,ls=1.25)
xr=MX+cw2+0.5
card(s,xr,y,cw2,hh,accent=True)
txt(s,xr+0.3,y+0.25,cw2-0.6,0.7,[[("113 m² ",BRANCO,True),("úteis",ICE2,True)]],26,BRANCO,bold=True,font=HEAD)
txt(s,xr+0.3,y+0.95,cw2-0.6,1.4,[[("A régua de VENDA. ",BRANCO,True),("É o piso interno do apartamento — a área do anúncio, dos portais e da conversa com o comprador. É com ela que nos posicionamos contra a concorrência.",ICE2,False)]],12.5,ICE2,ls=1.25)
sourcebox(s,MX,4.85,CW,1.0,[[("O fator do prédio é 1,88",TINTA,True),(" (213 ÷ 113) — alto e favorável: muita área comum e garagem por unidade, coerente com lazer completo rateado entre só 20 apartamentos. Na visita, antecipamos que \"213 m²\" da matrícula é a área total — evita ruído e transmite domínio.",CINZA700,False)]],size=12)
txt(s,MX,6.0,CW,0.6,"Misturar as duas réguas é erro de unidade — foi o que derrubou uma avaliação anterior deste imóvel em 47%. O nosso cálculo é cadastral contra cadastral: a única comparação tecnicamente válida.",11.5,CINZA500,ls=1.25)
foot(s,7)

# S8 METODOLOGIA
s=newslide(); kicker_title(s,"Metodologia · evidência rastreável, não opinião","Como chegamos ao valor")
cw3=(CW-0.6)/3; y=2.1; hh=2.15
stat(s,MX,y,cw3,hh,"631","","Apartamentos vendidos (ITBI/PMSP) num raio de 1.000 m, filtrados de 4.986 transações reais 2024–2026")
stat(s,MX+cw3+0.3,y,cw3,hh,"53","","Vendas do perfil exato do alvo — área cadastral 190–240 m² e prédios de 1985–2005 (o Manuela é de 1994)")
stat(s,MX+2*(cw3+0.3),y,cw3,hh,"IPTU","","Carnê do próprio edifício ancora a área oficial (213 m²) e o piso de valor (venal R$ 1,10 mi)",accent=True)
txt(s,MX,4.55,CW,0.8,[[("Cada comparável é uma ",CINZA700,False),("transação registrada em cartório",TINTA,True),(", rastreável pelo cadastro municipal (SQL). Preço de anúncio entra só como leitura de concorrência — ",CINZA700,False),("nunca como verdade de valor",TINTA,True),(".",CINZA700,False)]],15,CINZA700,ls=1.35)
sourcebox(s,MX,5.5,CW,1.1,[[("A idade do prédio é o fator que mais separa preço: ",TINTA,True),("no mesmo recorte de área, imóveis de todas as épocas têm mediana de R$ 1,95 mi — mas o estoque de 1985–2005 fica em R$ 1,31 mi. Comparar o Ed. Manuela com prédio novo inflaria o valor e travaria a venda.",CINZA700,False)]],size=12)
foot(s,8)

# S9 MAPA
s=newslide(); kicker_title(s,"Geografia · raio de 1.000 m a partir do imóvel","Os comparáveis no entorno imediato",tsize=27)
mapa = os.path.join(HERE, "ACM_alvorada_mapa.png")
if os.path.exists(mapa):
    s.shapes.add_picture(mapa, Inches(MX), Inches(2.0), height=Inches(4.7))
bullets(s,7.15,2.25,5.4,[
    [("Vermelho: ",RGBColor(0xE2,0x23,0x1A),True),("o imóvel-alvo (Rua Alvorada, 116 — apto 52).",CINZA700,False)],
    [("Verde: ",OK,True),("o comparável interno — AP 51, no próprio Ed. Manuela.",CINZA700,False)],
    [("Dourado (1–3): ",OURO,True),("os 3 comparáveis de máxima aderência.",CINZA700,False)],
    [("Laranja (4–5): ",LAR,True),("reforço da amostra.",CINZA700,False)],
    [("Azul: ",RGBColor(0x1E,0x5A,0xA8),True),("demais vendas do Top 20 no raio.",CINZA700,False)],
],size=13,gap=8)
sourcebox(s,7.15,5.35,5.4,1.15,"Microrregião de valorização homogênea na Vila Olímpia — 62 vendas reais do recorte, com o comparável de ouro dentro do próprio prédio.",size=11.5)
foot(s,9)

# S10 AS 3 EVIDÊNCIAS
s=newslide(); kicker_title(s,"A construção do número · três evidências independentes","Três evidências convergem para a mesma faixa",tsize=28)
cw3=(CW-0.6)/3; y=2.05; hh=3.3
card(s,MX,y,cw3,hh,accent=True)
txt(s,MX+0.25,y+0.2,cw3-0.5,0.6,"1 · O vizinho de andar",17,BRANCO,bold=True,font=HEAD)
txt(s,MX+0.25,y+0.8,cw3-0.5,2.3,[[("O apto 51 — mesmo prédio, mesmo andar, mesma área — vendido por ",ICE2,False),("R$ 1.325.000",BRANCO,True),(" (set/2023). Corrigido pelo micro-mercado: ",ICE2,False),("R$ 1,31–1,51 mi",BRANCO,True),(" hoje. Não existe comparável melhor em uma ACM.",ICE2,False)]],12,ICE2,ls=1.25)
x2=MX+cw3+0.3
card(s,x2,y,cw3,hh)
txt(s,x2+0.25,y+0.2,cw3-0.5,0.6,"2 · O prédio gêmeo",17,AZUL,bold=True,font=HEAD)
txt(s,x2+0.25,y+0.8,cw3-0.5,2.3,[[("Vahia de Abreu, 383 — também de 1994, a 390 m — teve ",CINZA700,False),("10 vendas",TINTA,True),(" em 2024–26, entre ",CINZA700,False),("R$ 1,00 e 1,24 mi",TINTA,True),(". O Manuela vale mais: é boutique, com 2 vagas e lazer.",CINZA700,False)]],12,CINZA700,ls=1.25)
x3=MX+2*(cw3+0.3)
card(s,x3,y,cw3,hh)
txt(s,x3+0.25,y+0.2,cw3-0.5,0.6,"3 · O mercado da região",17,AZUL,bold=True,font=HEAD)
txt(s,x3+0.25,y+0.8,cw3-0.5,2.3,[[("53 vendas reais",TINTA,True),(" do mesmo porte e da mesma safra no raio de 1.000 m — mediana de ",CINZA700,False),("R$ 1.312.500",TINTA,True),(", coincidindo com o piso da evidência nº 1.",CINZA700,False)]],12,CINZA700,ls=1.25)
sp=rect(s,MX,5.65,CW,0.85,AZUL,shape=MSO_SHAPE.ROUNDED_RECTANGLE); sp.adjustments[0]=0.10
txt(s,MX+0.3,5.73,CW-0.6,0.7,[[("Convergência: ",BRANCO,True),("o comparável interno corrigido (R$ 1,31–1,51 mi) e a mediana da mesma época (R$ 1,31 mi) apontam para a mesma região de valor — alta confiança na faixa.",ICE2,False)]],13.5,ICE2,anchor=MSO_ANCHOR.MIDDLE,ls=1.2)
foot(s,10)

# S11 COMPARÁVEL DE OURO (AP 51)
s=newslide(); kicker_title(s,"O comparável de ouro · vendas no próprio Ed. Manuela","A régua é uma venda no seu prédio, no seu andar",tsize=27)
table(s, MX, 2.05, CW,
    ["Unidade","Data","Área cad.","Valor","Natureza","Uso na ACM"],
    [
        ["AP 51 — 5º andar","set/2023","213 m²","R$ 1.325.000","Compra e venda","COMPARÁVEL DE OURO"],
        ["AP 61","set/2024","213 m²","R$ 800.000","Compra e venda","outlier — abaixo do venal"],
        ["AP 102","mar/2024","316 m²","R$ 1.153.774","Arrematação em leilão","excluído (não é mercado)"],
        ["AP 12","dez/2024","213 m²","R$ 200.000","Integralização de capital","excluído (não é venda)"],
    ],
    colw=[2.2,1.2,1.3,1.9,2.6,2.9], hl_row=0, fsize=11.5, rowh=0.42)
txt(s,MX,4.6,CW,0.95,[[("O AP 51 é o ",CINZA700,False),("vizinho de porta do 52, no mesmo 5º andar",TINTA,True),(": mesma área, mesma safra, mesmo condomínio, mesmo lazer. Corrigido pelo índice do micro-mercado, indica ",CINZA700,False),("R$ 1.308.000 (base 2025) a R$ 1.506.000 (base 2026 parcial)",TINTA,True),(" para o alvo hoje — antes dos diferenciais de conservação.",CINZA700,False)]],13,CINZA700,ls=1.3)
sourcebox(s,MX,5.7,CW,0.95,[[("Sobre o AP 61 (R$ 800 mil): ",TINTA,True),("transação isolada abaixo do próprio valor venal do prédio (R$ 1,10 mi) — padrão de venda atípica. Se um comprador citar, a resposta está pronta: contra ela há uma venda de mercado no mesmo andar a R$ 1,325 mi.",CINZA700,False)]],size=11.5)
foot(s,11)

# S12 TOP 5 ADERÊNCIA
s=newslide(); kicker_title(s,"Comparáveis de maior aderência · vendas registradas (ITBI)","Os 5 imóveis mais parecidos, já vendidos",tsize=28)
table(s, MX, 2.05, CW,
    ["#","Endereço (vendido)","Área cad.","Ano","Distância","Venda"],
    [
        ["1","R. Prof. Vahia de Abreu, 383","203 m²","1994","390 m","R$ 1.125.000"],
        ["2","R. Prof. Vahia de Abreu, 383","203 m²","1994","390 m","R$ 1.150.000"],
        ["3","Av. dos Eucaliptos, 113","214 m²","1996","789 m","R$ 1.300.000"],
        ["4","Av. dos Eucaliptos, 113","216 m²","1996","789 m","R$ 2.000.000"],
        ["5","Av. dos Eucaliptos, 155","205 m²","1995","829 m","R$ 1.150.000"],
        ["","Mediana dos 20 comparáveis de maior aderência (12 prédios)","","","","R$ 1.300.000"],
    ],
    colw=[0.7,4.6,1.4,1.1,1.5,2.8], avg_row=5, fsize=11.5, rowh=0.4)
txt(s,MX,5.1,CW,0.7,[[("Sinal de amostra estável: ",VERM,True),("o Top 20 e o recorte completo (62 vendas) dão o mesmo número — R$ 1.300.000. O Top 5 fica abaixo por concentrar no prédio gêmeo, de alta rotatividade e patamar inferior ao Manuela.",CINZA700,False)]],12.5,CINZA700,ls=1.25)
txt(s,MX,5.9,CW,0.5,"Fonte: ITBI/PMSP 2024–2026, rastreável por SQL cadastral. Aderência: 40% área · 25% época · 25% distância · 10% recência, com diversidade de prédios.",10.5,CINZA500,ls=1.2)
foot(s,12)

# S13 MERCADO FECHA ABAIXO + POSICIONAMENTO
s=newslide(); kicker_title(s,"A tese central · por que precificar certo protege o seu valor","O mercado fecha 25–30% abaixo do anúncio",tsize=27)
txt(s,MX,1.9,CW,0.6,[[("Na régua que o comprador usa (R$/m² útil), a faixa sugerida posiciona o seu apartamento ",CINZA700,False),("no meio exato do mercado ofertado",TINTA,True),(" — logo abaixo do concorrente direto da mesma rua:",CINZA700,False)]],13,CINZA700,ls=1.25)
table(s, MX, 2.6, CW,
    ["Concorrente (ofertas ativas)","Área útil","Pedido","R$/m² útil","vs. nosso anúncio"],
    [
        ["R. do Rocio","112 m²","R$ 2.499.000","R$ 22.313","muito acima"],
        ["R. Alvorada","109 m²","R$ 2.200.000","R$ 20.183","muito acima"],
        ["R. Alvorada","109 m²","R$ 1.700.000","R$ 15.596","logo acima do nosso teto"],
        ["SEU IMÓVEL (sugerido)","113 m²","R$ 1,55–1,70 mi","R$ 13.717–15.044","—"],
        ["R. Casa do Ator","115 m²","R$ 1.280.000","R$ 11.130","abaixo"],
        ["R. Alvorada (1 vaga)","110 m²","R$ 1.170.000","R$ 10.636","abaixo"],
    ],
    colw=[3.6,1.3,2.1,2.3,2.8], hl_row=3, fsize=11, rowh=0.38)
sourcebox(s,MX,5.5,CW,1.15,[[("Ofertas do entorno pedem mediana de ~R$ 1,70 mi; os fechamentos reais equivalentes ficam em R$ 1,20–1,65 mi (mediana R$ 1,31 mi) — ",CINZA700,False),("gap de 25–30%",TINTA,True),(". Anunciar acima do que o mercado paga gera menos visitas, mais prazo e maior perda ao final. ",CINZA700,False),("Precificar certo desde o início protege o seu patrimônio.",TINTA,True)]],size=11.5)
foot(s,13)

# S14 RECOMENDAÇÃO (HERO)
s=newslide(dark=True)
txt(s,MX,0.7,CW,0.4,"SUGESTÃO RE/MAX GALERIA · RECOMENDAÇÃO DE PRECIFICAÇÃO",13,ICE,bold=True)
txt(s,MX,1.25,CW,0.35,"Anúncio recomendado",13,ICE2)
txt(s,MX,1.6,7.5,1.1,[[("R$ ",ICE,True),("1,55–1,70",BRANCO,True),(" mi",ICE,True)]],54,BRANCO,bold=True,font=HEAD)
txt(s,7.6,1.25,CW-7.0,0.35,"Fechamento esperado",13,ICE2)
txt(s,7.6,1.6,5.5,1.0,[[("R$ ",ICE,True),("1,35–1,50",BRANCO,True),(" mi",ICE,True)]],40,BRANCO,bold=True,font=HEAD)
txt(s,MX,3.0,CW,0.35,"A RÉGUA COMPLETA DE VALOR (R$ MILHÕES)",12,ICE,bold=True)
cw5=(CW-4*0.22)/5; y=3.45; hh=1.35
labels=[("R$ 1,70","Anúncio — teto",False),("R$ 1,55","Anúncio — piso",False),("R$ 1,35–1,50","Fechamento esperado",True),("R$ 1,25","Piso de negociação",False),("R$ 1,10","Venal oficial (IPTU)",False)]
for i,(v,l,hl) in enumerate(labels):
    x=MX+i*(cw5+0.22)
    sp=rect(s,x,y,cw5,hh, VERM if hl else AZUL700, shape=MSO_SHAPE.ROUNDED_RECTANGLE); sp.adjustments[0]=0.10
    txt(s,x+0.1,y+0.22,cw5-0.2,0.5,v,15 if hl else 17,BRANCO,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
    txt(s,x+0.1,y+0.78,cw5-0.2,0.5,l,10.5,RGBColor(0xff,0xd9,0xdf) if hl else ICE2,align=PP_ALIGN.CENTER,ls=1.1)
rect(s,MX,5.15,0.06,1.5,VERM)
txt(s,MX+0.25,5.2,CW-0.5,1.5,[
    [("A régua é o seu próprio prédio: ",BRANCO,True),("a venda do AP 51 corrigida (R$ 1,31–1,51 mi) mais os diferenciais do 52 — conservação boa (+3 a +6%) e ventilação cruzada (+2 a +4%).",ICE2,False)],
    [("Anunciar a R$ 1,55–1,70 mi posiciona na faixa do concorrente de R$ 1,70 mi da mesma rua, com ~15% de reserva. ",ICE2,False),("Não descer abaixo de R$ 1,25 mi",BRANCO,True),(" — defendido pelo venal oficial e pela venda do vizinho.",ICE2,False)],
],12.5,ICE2,ls=1.3,space=6)
foot(s,14,dark=True)

# S15 PRÓXIMOS PASSOS
s=newslide(); kicker_title(s,"Estratégia de venda · próximos passos","Do contrato à venda")
steps=[
    ("1","Assinatura do contrato de representação exclusiva.",False),
    ("2","Análise jurídica da documentação (due diligence) — carnê de IPTU do apto 52 e averbação das 2 vagas na matrícula.",False),
    ("3","Montagem do plano de marketing e produção do material profissional (fotos, filme e tour virtual).",False),
    ("4","Implementação e início dos trabalhos — anúncio entre R$ 1,55 e 1,70 mi, com revisão em ciclos de 45–60 dias.",False),
    ("5","VENDA — fechamento entre R$ 1,35 e R$ 1,50 milhão (piso R$ 1,25 mi).",True),
]
y=2.15
for n,t,win in steps:
    c = VERM if win else AZUL
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MX), Inches(y), Inches(0.5), Inches(0.5))
    sp.fill.solid(); sp.fill.fore_color.rgb=c; sp.line.fill.background(); _noshadow(sp)
    tf=sp.text_frame; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=n; r.font.size=Pt(15); r.font.bold=True; r.font.name=HEAD; r.font.color.rgb=BRANCO
    txt(s,MX+0.75,y+0.02,CW-1.0,0.75,t,14.5 if not win else 16,(VERM if win else CINZA700),bold=win,ls=1.2)
    y+=0.92
foot(s,15)

# S16 ENCERRAMENTO
s=prs.slides.add_slide(BLANK); bg(s,AZUL900); brandbar(s)
txt(s,MX,1.35,CW,0.4,"RE/MAX GALERIA · MOEMA",14,ICE,bold=True)
txt(s,MX,1.85,9.5,1.1,"Obrigada, Dorothy.",44,BRANCO,bold=True,font=HEAD)
txt(s,MX,2.95,8.4,0.9,[[("Se deseja vender seu imóvel com ",ICE2,False),("gestão profissional focada em resultados",BRANCO,True),(", fale comigo.",ICE2,False)]],18,ICE2,ls=1.35)
rect(s,MX,4.0,0.06,1.15,VERM)
txt(s,MX+0.25,4.05,8.6,1.1,"RE/MAX Galeria — conectando pessoas aos melhores negócios imobiliários por um modelo de gestão eficiente, transparente e orientado por resultados.",13.5,ICE2,italic=True,ls=1.35)
txt(s,7.7,5.1,5.0,1.8,[
    [("Luciana Borba",BRANCO,True)],
    [("Consultora Imobiliária · RE/MAX Galeria · (11) 99995-2014",ICE2,False)],
    [("lucianaborba@remax.com.br",ICE2,False)],
    [("CRECI: 045063-J",ICE,False)],
],12.5,ICE2,align=PP_ALIGN.RIGHT,ls=1.3,space=2)
bug(s,MX,5.9)

out = os.path.join(HERE, "PITCH_Captacao_Alvorada116_Dorothy_RE-MAX_v1.pptx")
prs.save(out)
print("OK:", out)
